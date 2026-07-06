# agy/utils/file_readers.py

"""Helper routines to load different file formats as plain text."""

from __future__ import annotations

import csv
from collections.abc import Callable, Iterable
from pathlib import Path

import pdfplumber
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def find_file_in_standard_dirs(
    filename: str, search_dirs: list[str] | None = None
) -> Path:
    """
    Find a file using fallback strategy:
    1. First try: in standard directories relative to cwd (., prompts, data, objects)
    2. Second try: in standard directories relative to project root
    3. If not found: raise FileNotFoundError

    Args:
        filename: Name of file to find (can include subdirectories or absolute path)
        search_dirs: Optional custom search directories (default: FILE_SEARCH_ORDER from config)

    Returns:
        Resolved Path object

    Raises:
        FileNotFoundError: If file not found in any search location
    """
    from agy.config import FILE_SEARCH_ORDER, get_project_root

    # Handle absolute paths - check if they exist directly
    file_path = Path(filename)
    if file_path.is_absolute():
        if file_path.exists():
            return file_path.resolve()
        raise FileNotFoundError(f"File not found: {file_path}")

    # Use provided search_dirs or default from config
    dirs_to_search = search_dirs or FILE_SEARCH_ORDER
    searched: list[str] = []

    # ============================================================================
    # STEP 1: Try in standard directories relative to cwd
    # ============================================================================
    base_path_cwd = Path.cwd()

    for search_dir in dirs_to_search:
        if search_dir == ".":
            candidate = base_path_cwd / filename
        else:
            candidate = base_path_cwd / search_dir / filename
        searched.append(str(candidate))
        if candidate.exists():
            return candidate.resolve()

    # ============================================================================
    # STEP 2: Try in standard directories relative to project root
    # ============================================================================
    project_root = get_project_root()

    # Only try project root if it's different from cwd (avoid duplicate search)
    if project_root.resolve() != base_path_cwd.resolve():
        for search_dir in dirs_to_search:
            if search_dir == ".":
                candidate = project_root / filename
            else:
                candidate = project_root / search_dir / filename
            searched.append(str(candidate))
            if candidate.exists():
                return candidate.resolve()

    # ============================================================================
    # STEP 3: Not found - raise error with helpful message
    # ============================================================================
    raise FileNotFoundError(
        f"File '{filename}' not found in any of these locations:\n"
        + "\n".join(f"  - {loc}" for loc in searched)
    )


def _read_txt(path: Path) -> str:
    """Read txt.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    return path.read_text(encoding="utf-8")


def _read_csv(path: Path) -> str:
    """Read csv.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    rows: list[str] = []
    with path.open("r", encoding="utf-8", newline="") as file:
        reader = csv.reader(file)
        for row in reader:
            rows.append(", ".join(row))
    return "\n".join(rows)


def _read_pdf(path: Path) -> str:
    """Read pdf.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    texts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            texts.append(page.extract_text() or "")
    return "\n".join(texts).strip()


def _read_docx(path: Path) -> str:
    """Read docx.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    document = Document(str(path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _read_pptx(path: Path) -> str:
    """Read pptx.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    presentation = Presentation(str(path))
    slide_texts: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
        slide_texts.append("\n".join(parts))
    return "\n\n".join(slide_texts)


def _read_xlsx(path: Path) -> str:
    """Read xlsx.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheet_texts: list[str] = []
    for sheet in workbook:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [""]
            if row:
                cells = ["" if value is None else str(value) for value in row]
            rows.append(", ".join(cells))
        sheet_texts.append(f"# Sheet: {sheet.title}\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(sheet_texts)


def _read_html(path: Path) -> str:
    """Read html.

    Args:
        path: path.

    Returns:
        str: Operation result.
    """
    content = path.read_text(encoding="utf-8")
    soup = BeautifulSoup(content, "html.parser")
    return str(soup.get_text(separator="\n"))


_READERS: dict[str, Callable[[Path], str]] = {
    ".txt": _read_txt,
    ".md": _read_txt,
    ".csv": _read_csv,
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".xlsx": _read_xlsx,
    ".htm": _read_html,
    ".html": _read_html,
}


def load_files_as_text(paths: Iterable[Path]) -> str:
    """Load multiple files and concatenate their contents with headers."""
    sections: list[str] = []
    for path in paths:
        suffix = path.suffix.lower()
        if suffix not in _READERS:
            raise ValueError(f"Unsupported file type: {path.suffix}")
        reader = _READERS[suffix]
        text = reader(path)
        sections.append(f"*** {path.name} ***\n{text}\n")
    return "\n".join(sections).strip()
